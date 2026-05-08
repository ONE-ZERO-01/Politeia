#!/usr/bin/env python3
"""Generate a project presentation for Politeia."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID    = RGBColor(0x16, 0x21, 0x3E)
BG_CARD   = RGBColor(0x1F, 0x2B, 0x4D)
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2   = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT3   = RGBColor(0x4E, 0xCD, 0xC4)
ACCENT4   = RGBColor(0xFF, 0xE6, 0x6D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
DIM       = RGBColor(0x99, 0x99, 0x99)
ORANGE    = RGBColor(0xFF, 0xA5, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# ── helpers ──
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _tf(shape):
    return shape.text_frame


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_paragraph(para, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    para.text = text
    para.alignment = align
    for run in para.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name


def add_paragraph(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2), font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return p


def first_paragraph(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(2)
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return p


def add_card(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def title_slide_header(slide, title, subtitle=None):
    add_accent_line(slide, Inches(0.8), Inches(0.5), Inches(3), ACCENT)
    tb = add_text_box(slide, Inches(0.8), Inches(0.65), Inches(11.5), Inches(1.0))
    tf = _tf(tb)
    tf.word_wrap = True
    first_paragraph(tf, title, size=32, bold=True, color=WHITE)
    if subtitle:
        add_paragraph(tf, subtitle, size=16, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, BG_DARK)

add_accent_line(s, Inches(1.5), Inches(2.2), Inches(4), ACCENT)

tb = add_text_box(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "Politeia", size=54, bold=True, color=WHITE)

tb = add_text_box(s, Inches(1.5), Inches(3.7), Inches(10), Inches(1.0))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "基于 Langevin 动力学的人类文明演化模拟器", size=28, color=ACCENT)

tb = add_text_box(s, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "Langevin-Jump Diffusion · C++20 / MPI+OpenMP · 31 Phases Completed", size=16, color=DIM)

tb = add_text_box(s, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "2026.05", size=14, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Core Idea
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "核心思想", "从物理学到文明演化")

tb = add_text_box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.2))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf,
    "用 Langevin 动力学框架模拟人类社会活动：将每个人抽象为\u201c粒子\u201d，"
    "在真实世界地形构成的势能面上，受确定性力、耗散力和随机力的共同驱动而演化。",
    size=18, color=LIGHT)

mapping = [
    ("粒子",         "个体人类",     "每个 agent 代表一个人"),
    ("位置 & 速度",  "地理位置 + 迁徙", "结合真实 DEM 地形数据"),
    ("势能函数 V(x)","地形·资源·人际力","核心难点：合理的势函数设计"),
    ("温度 T",       "社会\u201c活跃度\u201d",  "控制随机涨落，高温 = 乱世"),
    ("摩擦 −γv",     "制度腐化·遗忘",  "耗散项，使系统不可逆"),
    ("随机力 ξ(t)",  "日常随机事件",   "高斯白噪声，微观社会涨落"),
    ("Poisson 跳跃", "技术突破·大事件", "非连续跃迁 (ε / w / 瘟疫)"),
    ("相变",         "制度突变 (革命)", "最有价值的观测目标"),
    ("凝聚 / 成核",  "聚集为部落·国家", "涌现行为的核心"),
]

card_left = Inches(0.8)
card_top = Inches(3.2)
col_w = [Inches(2.0), Inches(2.5), Inches(7.0)]

header_card = add_card(s, card_left, card_top, Inches(11.5), Inches(0.45), fill_color=RGBColor(0x0A, 0x3D, 0x62))
htb = add_text_box(s, card_left + Inches(0.15), card_top + Inches(0.05), Inches(11.2), Inches(0.35))
htf = _tf(htb)
first_paragraph(htf, "分子动力学概念              社会映射                    说明", size=13, bold=True, color=ACCENT)

for i, (md_concept, soc, note) in enumerate(mapping):
    y = card_top + Inches(0.5) + Inches(i * 0.4)
    bg_c = BG_CARD if i % 2 == 0 else RGBColor(0x17, 0x24, 0x42)
    add_card(s, card_left, y, Inches(11.5), Inches(0.38), fill_color=bg_c)
    row_tb = add_text_box(s, card_left + Inches(0.15), y + Inches(0.04), Inches(11.2), Inches(0.32))
    row_tf = _tf(row_tb)
    first_paragraph(row_tf, f"{md_concept:<20s}{soc:<22s}{note}", size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Agent Definition
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "Agent 定义：个体状态向量", "Z_i(t) = (x, p, w, c⃗, ε)")

eq_tb = add_text_box(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.6))
eq_tf = _tf(eq_tb)
first_paragraph(eq_tf, "四个内部自由度对应人的四个不可约维度：在 (x)、有 (w)、知 (c⃗)、能 (ε)", size=17, color=ACCENT3)

dofs = [
    ("x_i", "地理位置 (2D)", "物质存在", "粒子位置", ACCENT),
    ("p_i", "动量 (2D)", "物质存在", "粒子动量", ACCENT),
    ("w_i", "财富 / 资源存量", "生存需求", "势能（可释放存量）", ACCENT2),
    ("c⃗_i", "文化向量 (d 维)", "认知 + 认同", "内部自旋 / 取向", ACCENT3),
    ("ε_i", "能量利用能力", "驾驭自然", "耦合常数", ACCENT4),
]

for i, (sym, meaning, phil, phys, clr) in enumerate(dofs):
    x_off = Inches(0.8)
    y_off = Inches(2.7) + Inches(i * 0.85)
    card = add_card(s, x_off, y_off, Inches(11.5), Inches(0.75), border_color=clr)

    tb1 = add_text_box(s, x_off + Inches(0.2), y_off + Inches(0.08), Inches(1.5), Inches(0.35))
    first_paragraph(_tf(tb1), sym, size=22, bold=True, color=clr)
    tb2 = add_text_box(s, x_off + Inches(1.8), y_off + Inches(0.08), Inches(3.5), Inches(0.35))
    first_paragraph(_tf(tb2), meaning, size=16, color=WHITE)
    tb3 = add_text_box(s, x_off + Inches(5.5), y_off + Inches(0.08), Inches(2.5), Inches(0.35))
    first_paragraph(_tf(tb3), phil, size=14, color=DIM)
    tb4 = add_text_box(s, x_off + Inches(8.2), y_off + Inches(0.08), Inches(3.0), Inches(0.35))
    first_paragraph(_tf(tb4), phys, size=14, color=DIM)

    if sym == "c⃗_i":
        sub_tb = add_text_box(s, x_off + Inches(1.8), y_off + Inches(0.4), Inches(9.0), Inches(0.3))
        first_paragraph(_tf(sub_tb), "模 |c⃗| = 知识量 → 影响 ε 跃迁概率    方向 ĉ = 文化取向 → 定义群体认同", size=11, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Physics Framework: Langevin Equation
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "物理框架：Langevin-跳跃扩散方程", "m·q̈ = −∇V − γ·q̇ + ξ(t) + J(t)")

terms = [
    ("−∇V(q)", "确定性力", "地形势 + 人际 LJ 势 → 聚集、迁徙", "force/", ACCENT),
    ("−γ·q̇", "耗散力", "制度腐化、知识遗忘 → 系统不可逆", "integrator/", ACCENT2),
    ("ξ(t)", "高斯白噪声", "日常随机事件（布朗运动）", "integrator/", ACCENT3),
    ("J(t)", "Poisson 跳跃", "技术突破、财富巨变、瘟疫", "interaction/", ACCENT4),
]

for i, (sym, name, desc, mod, clr) in enumerate(terms):
    x_off = Inches(0.8)
    y_off = Inches(2.0) + Inches(i * 1.15)
    card = add_card(s, x_off, y_off, Inches(11.5), Inches(1.0), border_color=clr)

    tb_sym = add_text_box(s, x_off + Inches(0.3), y_off + Inches(0.1), Inches(2.0), Inches(0.4))
    first_paragraph(_tf(tb_sym), sym, size=22, bold=True, color=clr)
    tb_name = add_text_box(s, x_off + Inches(2.5), y_off + Inches(0.1), Inches(2.5), Inches(0.4))
    first_paragraph(_tf(tb_name), name, size=18, bold=True, color=WHITE)
    tb_desc = add_text_box(s, x_off + Inches(2.5), y_off + Inches(0.5), Inches(6.5), Inches(0.4))
    first_paragraph(_tf(tb_desc), desc, size=14, color=LIGHT)
    tb_mod = add_text_box(s, x_off + Inches(9.5), y_off + Inches(0.25), Inches(1.8), Inches(0.4))
    first_paragraph(_tf(tb_mod), mod, size=13, color=DIM, align=PP_ALIGN.RIGHT)

bbk_tb = add_text_box(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6))
bbk_tf = _tf(bbk_tb)
first_paragraph(bbk_tf, "积分方法：BBK (Brünger-Brooks-Karplus)  |  涨落-耗散关系：σ² = 2γmkT  |  γ=0, T=0 → 退化为 Velocity-Verlet（能量守恒）", size=13, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Interaction Mechanisms
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "核心交互机制", "对称规则 → 不对称结果（涌现方法论）")

mechanisms = [
    ("资源交换", "Δw = η·(A_i−A_j)/(A_i+A_j)·min(w_i,w_j)", "规则对称，结果不对称；不平等自发涌现", ACCENT),
    ("文化传播", "Δc⃗ = rate·exp(−|Δc⃗|²/2σ²)·(c⃗_j−c⃗_i)·dt", "文化距离小 → 趋同；距离大 → 几乎不互动", ACCENT3),
    ("技术演化 ε", "漂移 + 接触传播 + Poisson 跳跃", "缓慢积累 + 技术突破（火、冶铁、蒸汽机）", ACCENT4),
    ("忠诚度 L", "dL/dt = α·保护 − β·税收 − γ·|Δc⃗| + η", "叛乱 / 投靠 / 征服 → 政体兴衰", ACCENT2),
    ("瘟疫", "免疫向量 d⃗ + SIR 传播 + 差异化打击", "大陆文明免疫谱广 → 跨群体接触灾难", ORANGE),
]

for i, (name, formula, desc, clr) in enumerate(mechanisms):
    y = Inches(2.0) + Inches(i * 1.0)
    card = add_card(s, Inches(0.8), y, Inches(11.5), Inches(0.9), border_color=clr)

    tb_n = add_text_box(s, Inches(1.1), y + Inches(0.08), Inches(2.0), Inches(0.35))
    first_paragraph(_tf(tb_n), name, size=18, bold=True, color=clr)
    tb_f = add_text_box(s, Inches(3.2), y + Inches(0.08), Inches(8.8), Inches(0.35))
    first_paragraph(_tf(tb_f), formula, size=14, bold=True, color=WHITE, font_name="Consolas")
    tb_d = add_text_box(s, Inches(3.2), y + Inches(0.48), Inches(8.8), Inches(0.35))
    first_paragraph(_tf(tb_d), desc, size=13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Population Dynamics
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "人口动力学", "繁殖·死亡·代际传承·婚配涌现")

# Left column - reproduction
card_l = add_card(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8))
tb = add_text_box(s, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.4))
first_paragraph(_tf(tb), "繁殖模型", size=20, bold=True, color=ACCENT)

repro_items = [
    "生育能力 φ(a) = Beta 分布钟形曲线",
    "五条件：距离 + 存活年龄 + 冷却期 + 文化兼容 + 财富门槛",
    "性别非对称：♀ 冷却期 2.75年/胎  ♂ 极短",
    "婚配制度自发涌现：贫穷→一夫一妻  富裕→一夫多妻",
    "后代继承：位置中点 + 文化均值 + ε 取较高者",
    "马尔萨斯反馈：ρ/K 抑制生育率和人均产出",
]
for j, item in enumerate(repro_items):
    tb = add_text_box(s, Inches(1.3), Inches(2.65) + Inches(j * 0.55), Inches(4.8), Inches(0.5))
    ttf = _tf(tb)
    tf_w = ttf
    tf_w.word_wrap = True
    first_paragraph(tf_w, f"▸ {item}", size=13, color=LIGHT)

# Right column - death
card_r = add_card(s, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8))
tb = add_text_box(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.4))
first_paragraph(_tf(tb), "四重死亡机制", size=20, bold=True, color=ACCENT2)

death_items = [
    ("衰老", "P ∝ exp(β·age)  [Gompertz 法则]", ACCENT2),
    ("饥饿", "w < threshold → 死亡率急剧增加", ACCENT2),
    ("意外", "P = λ（常数）——命运，不分贵贱", ACCENT2),
    ("瘟疫", "P = P₀·(1−d^k)·(1−ε/ε_med)", ACCENT2),
]
for j, (nm, fm, clr) in enumerate(death_items):
    y = Inches(2.7) + Inches(j * 0.7)
    tb_nm = add_text_box(s, Inches(7.3), y, Inches(1.5), Inches(0.3))
    first_paragraph(_tf(tb_nm), nm, size=16, bold=True, color=clr)
    tb_fm = add_text_box(s, Inches(8.8), y, Inches(3.5), Inches(0.3))
    first_paragraph(_tf(tb_fm), fm, size=13, color=LIGHT, font_name="Consolas")

# Lifespan emergence
tb = add_text_box(s, Inches(7.1), Inches(5.3), Inches(5.2), Inches(1.2))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "寿命作为派生量自然涌现", size=15, bold=True, color=ACCENT4)
add_paragraph(tf, "λ_i = h(w_i, ε_local)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "原始社会 ⟨a⟩≈35 → 农业革命 ⟨a⟩≈40 → 工业革命 ⟨a⟩≈75", size=12, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Hierarchy & Polity
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "层级、忠诚度与政体涌现", "从个体对称规则到帝国制度")

# Left: Loyalty system
card_l = add_card(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.8))
tb = add_text_box(s, Inches(1.1), Inches(2.1), Inches(5.4), Inches(0.4))
first_paragraph(_tf(tb), "忠诚度演化方程", size=20, bold=True, color=ACCENT)

loyalty_items = [
    ("保护收益 (+α)", "上级分配资源 → 忠诚增加", ACCENT3),
    ("税收消耗 (−β)", "过度征税 → 忠诚降低", ACCENT2),
    ("文化距离 (−γ)", "异族统治天然不稳定", ACCENT4),
    ("叛乱 L < 0.1", "断裂依附，成为独立根节点", ACCENT2),
    ("投靠 L < 0.2", "转移到更强的领主", ORANGE),
    ("征服", "Power_i > 1.5×Power_j → 吞并", ACCENT),
]
for j, (nm, desc, clr) in enumerate(loyalty_items):
    y = Inches(2.65) + Inches(j * 0.6)
    tb_nm = add_text_box(s, Inches(1.3), y, Inches(5.2), Inches(0.3))
    first_paragraph(_tf(tb_nm), f"▸ {nm}", size=14, bold=True, color=clr)
    tb_d = add_text_box(s, Inches(1.5), y + Inches(0.25), Inches(5.0), Inches(0.3))
    first_paragraph(_tf(tb_d), desc, size=12, color=DIM)

# Right: Emergence results
card_r = add_card(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.8))
tb = add_text_box(s, Inches(7.3), Inches(2.1), Inches(5.0), Inches(0.4))
first_paragraph(_tf(tb), "涌现验证结果 (1000 粒子)", size=18, bold=True, color=ACCENT3)

metrics = [
    ("H (层级深度)", "5—7 层", "从无层级到官僚帝国"),
    ("C (独立政体)", "110\u2014254 个", "自发形成多个\u201c国家\u201d"),
    ("F (最大帝国)", "3—10%", "最大帝国占总人口比"),
    ("Ψ (封建化)", "0.52—0.69", "中等封建结构"),
    ("Gini (财富)", "~0.89", "对称规则→高度不平等"),
    ("⟨L⟩ (忠诚度)", "0.78—0.94", "世袭使忠诚更稳定"),
]
for j, (name, val, desc) in enumerate(metrics):
    y = Inches(2.7) + Inches(j * 0.6)
    tb1 = add_text_box(s, Inches(7.3), y, Inches(2.2), Inches(0.3))
    first_paragraph(_tf(tb1), name, size=13, color=LIGHT)
    tb2 = add_text_box(s, Inches(9.5), y, Inches(1.5), Inches(0.3))
    first_paragraph(_tf(tb2), val, size=14, bold=True, color=ACCENT)
    tb3 = add_text_box(s, Inches(7.5), y + Inches(0.25), Inches(4.8), Inches(0.25))
    first_paragraph(_tf(tb3), desc, size=11, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Geography & Terrain Experiments
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "地理决定论验证：地形对比实验", "山脉屏障 → 碎裂化 ｜ 开放平原 → 集中化")

tb = add_text_box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "核心假说：地理约束驱动政治集中化。通过 6 种合成地形 × 相同参数的对比实验验证。", size=16, color=LIGHT)

terrain_data = [
    ("中国 (大平原)", "H=7", "15个政体", "HHI=0.088", "最深层级、最少政体、较高集中度", ACCENT2),
    ("欧洲 (山脉分割)", "H=5", "27个政体", "HHI=0.052", "最多政体、较浅层级、最低集中度", ACCENT3),
    ("河谷 (River)", "H=5", "—", "HHI=0.08", "最高 HHI 集中度", ACCENT),
    ("高斯单中心", "H=3", "30个政体", "—", "最碎片化，最浅层级", ACCENT4),
]

headers = ["地形类型", "层级深度 H", "政体数", "HHI 集中度", "关键发现"]
hdr_card = add_card(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.45), fill_color=RGBColor(0x0A, 0x3D, 0x62))
htb = add_text_box(s, Inches(1.0), Inches(2.85), Inches(11.2), Inches(0.35))
htf = _tf(htb)
first_paragraph(htf, f"{'地形类型':<18s}{'层级 H':<12s}{'政体数':<12s}{'HHI':<14s}{'关键发现'}", size=13, bold=True, color=ACCENT)

for i, (terrain, h, pol, hhi, finding, clr) in enumerate(terrain_data):
    y = Inches(3.3) + Inches(i * 0.55)
    bg_c = BG_CARD if i % 2 == 0 else RGBColor(0x17, 0x24, 0x42)
    add_card(s, Inches(0.8), y, Inches(11.5), Inches(0.48), fill_color=bg_c)
    row_tb = add_text_box(s, Inches(1.0), y + Inches(0.08), Inches(11.2), Inches(0.32))
    row_tf = _tf(row_tb)
    first_paragraph(row_tf, f"{terrain:<18s}{h:<12s}{pol:<12s}{hhi:<14s}{finding}", size=13, color=WHITE)

# conclusion box
con_card = add_card(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.2), border_color=ACCENT4)
tb = add_text_box(s, Inches(1.1), Inches(5.8), Inches(11.0), Inches(1.0))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "关键发现", size=16, bold=True, color=ACCENT4)
add_paragraph(tf, "• 山脉屏障阻止政体扩张 → 碎裂化；开放平原允许兼并 → 集中化", size=14, color=LIGHT)
add_paragraph(tf, "• Gini 不平等在所有地形中恒定 (~0.89) — 不平等是对称规则的内在结果，而非地理决定", size=14, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Technical Architecture
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "技术架构", "C++20 · SoA · MPI+OpenMP · Morton Z-curve")

# Tech stack cards
stack = [
    ("C++20", "语言标准", "std::ranges, concepts, modules-ready", ACCENT),
    ("SoA 数据布局", "内存优化", "连续数组存储，缓存友好，SIMD 向量化", ACCENT3),
    ("BBK 积分器", "时间步进", "Langevin 方程二阶积分，涨落-耗散自洽", ACCENT4),
    ("Cell List", "邻居搜索", "O(N) 桶排序，每步重建", ACCENT2),
]

for i, (name, cat, desc, clr) in enumerate(stack):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(2.0) + row * Inches(1.2)
    card = add_card(s, x, y, Inches(5.5), Inches(1.0), border_color=clr)
    tb_n = add_text_box(s, x + Inches(0.2), y + Inches(0.08), Inches(3.0), Inches(0.35))
    first_paragraph(_tf(tb_n), name, size=18, bold=True, color=clr)
    tb_c = add_text_box(s, x + Inches(3.3), y + Inches(0.08), Inches(2.0), Inches(0.35))
    first_paragraph(_tf(tb_c), cat, size=13, color=DIM, align=PP_ALIGN.RIGHT)
    tb_d = add_text_box(s, x + Inches(0.2), y + Inches(0.5), Inches(5.0), Inches(0.4))
    first_paragraph(_tf(tb_d), desc, size=13, color=LIGHT)

# Parallel architecture
par_card = add_card(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), border_color=ACCENT)
tb = add_text_box(s, Inches(1.1), Inches(4.6), Inches(11.0), Inches(0.4))
first_paragraph(_tf(tb), "并行架构：MPI + OpenMP 混合", size=18, bold=True, color=ACCENT)

par_items = [
    "MPI 区域分解 → Morton Z-curve 空间填充曲线 → 负载均衡效率 97%+",
    "OpenMP 线程级并行 → 力计算、积分器、资源动力学等热点路径",
    "粒子迁移：跨域 MPI 通信含 global_id / superior / loyalty",
    "弹性 Checkpoint/Restart：N 进程写 → M 进程读（如 4→400）",
    "Scaling 实测：Strong (N=8000, P=1→4) speedup 1.65x；Weak 负载均衡效率 97%",
]
for j, item in enumerate(par_items):
    tb_item = add_text_box(s, Inches(1.3), Inches(5.15) + Inches(j * 0.35), Inches(10.8), Inches(0.32))
    first_paragraph(_tf(tb_item), f"▸ {item}", size=13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Source Code Modules
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "源码模块总览", "src/ 目录结构与物理对应")

modules = [
    ("core/", "基础数据层", "ParticleData (SoA), Config, 常量, 类型", ACCENT),
    ("domain/", "空间管理", "CellList 邻居搜索 · SFC 分解 · Halo 交换", ACCENT3),
    ("force/", "力计算", "人际 LJ 势 · 地形外势 · TerrainGrid · RiverField", ACCENT2),
    ("integrator/", "时间积分", "BBK Langevin 积分器 · Poisson 跳跃 · RNG", ACCENT4),
    ("interaction/", "个体交互", "资源交换 · 文化传播 · 技术演化 · 忠诚度 · 征服", ORANGE),
    ("population/", "人口动力学", "繁殖 · 死亡 · 遗传 · 瘟疫 · 承载力", ACCENT2),
    ("analysis/", "分析观测", "序参量 · 网络分析 · 政体检测 · 性能监控", ACCENT3),
    ("io/", "输入输出", "CSV 写入 · IC 加载 · Checkpoint/Restart", DIM),
]

for i, (mod, role, desc, clr) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(2.0) + row * Inches(1.2)
    card = add_card(s, x, y, Inches(5.5), Inches(1.05), border_color=clr)
    tb_m = add_text_box(s, x + Inches(0.2), y + Inches(0.08), Inches(2.0), Inches(0.35))
    first_paragraph(_tf(tb_m), mod, size=18, bold=True, color=clr, font_name="Consolas")
    tb_r = add_text_box(s, x + Inches(2.3), y + Inches(0.08), Inches(3.0), Inches(0.35))
    first_paragraph(_tf(tb_r), role, size=14, color=WHITE)
    tb_d = add_text_box(s, x + Inches(0.2), y + Inches(0.5), Inches(5.0), Inches(0.45))
    dtf = _tf(tb_d)
    dtf.word_wrap = True
    first_paragraph(dtf, desc, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Development Progress
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "开发进度：31 个 Phase 全部完成", "~44 周 · 涌现优先原则")

phases_summary = [
    ("Phase 0-3", "基础骨架 + 空间动力学 + 资源交换 + Langevin", "✅"),
    ("Phase 4-6", "文化向量 + 人口动力学 + ε 与 Poisson 跳跃", "✅"),
    ("Phase 7-9", "MPI 并行 + 交互网络 + 分析 + 可视化 + 瘟疫", "✅"),
    ("Phase 10-12", "Morton 曲线 SFC + 真实地形 + MPI+OpenMP 混合", "✅"),
    ("Phase 13-16", "忠诚度 + 世袭 + 显式层级分析 + 可视化更新", "✅"),
    ("Phase 17-22", "Scaling 测试 + 数据增强 + 合成地形 + 政体检测 + 相变", "✅"),
    ("Phase 23-26", "承载力 + 地形对比 + 中国vs欧洲 + 可配置实验框架", "✅"),
    ("Phase 27-31", "世界地图 + 初始条件 + Checkpoint + HYDE 校准 + 河流场", "✅"),
]

for i, (phase, desc, status) in enumerate(phases_summary):
    y = Inches(2.0) + Inches(i * 0.6)
    bg_c = RGBColor(0x0A, 0x3D, 0x62) if i % 2 == 0 else BG_CARD
    add_card(s, Inches(0.8), y, Inches(11.5), Inches(0.52), fill_color=bg_c)
    tb_p = add_text_box(s, Inches(1.0), y + Inches(0.08), Inches(2.0), Inches(0.35))
    first_paragraph(_tf(tb_p), phase, size=15, bold=True, color=ACCENT)
    tb_d = add_text_box(s, Inches(3.2), y + Inches(0.08), Inches(8.0), Inches(0.35))
    first_paragraph(_tf(tb_d), desc, size=13, color=LIGHT)
    tb_s = add_text_box(s, Inches(11.5), y + Inches(0.08), Inches(0.7), Inches(0.35))
    first_paragraph(_tf(tb_s), status, size=16, color=ACCENT3, align=PP_ALIGN.CENTER)

# Stats bar
stats_bar = add_card(s, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.5), border_color=ACCENT)
stats = [
    ("110+", "单元测试"),
    ("31", "开发阶段"),
    ("~44 周", "开发周期"),
    ("10⁵", "粒子规模验证"),
]
for j, (num, label) in enumerate(stats):
    x_s = Inches(1.2) + Inches(j * 3.0)
    tb_n = add_text_box(s, x_s, Inches(6.92), Inches(1.0), Inches(0.25))
    first_paragraph(_tf(tb_n), num, size=18, bold=True, color=ACCENT)
    tb_l = add_text_box(s, x_s + Inches(1.1), Inches(6.95), Inches(1.6), Inches(0.25))
    first_paragraph(_tf(tb_l), label, size=13, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Key Scientific Findings
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "核心科学发现", "对称规则产生不对称结果")

findings = [
    ("不平等自发涌现",
     "Gini 系数从 ~0 自发上升至 ~0.89，财富分布从均匀分布变为 Pareto 幂律分布。"
     "对称的交换规则在 ε 差异放大下产生马太效应。",
     ACCENT),
    ("文化圈自发形成",
     "文化向量的局部同化 + 远距离排斥使文化边界从连续空间中自然涌现，"
     "无需预设离散标签。",
     ACCENT3),
    ("政治制度自发涌现",
     "从无层级的平等个体出发，依次涌现：Band → Tribe → Chiefdom → State → Empire。"
     "层级深度 H 从 0 增长到 5-7。",
     ACCENT4),
    ("地理决定政治结构",
     "中国式大平原 → 深层级、少政体、高集中度；欧洲式山脉分割 → 浅层级、多政体、低集中度。"
     "验证地理约束驱动政治集中化假说。",
     ACCENT2),
    ("婚配制度自发涌现",
     "从生物非对称性（妊娠冷却期）+ 财富分化 + 权力集中，"
     "一夫多妻、后宫制度无需预设即可自然生长出来。",
     ORANGE),
]

for i, (title, desc, clr) in enumerate(findings):
    y = Inches(1.9) + Inches(i * 1.05)
    card = add_card(s, Inches(0.8), y, Inches(11.5), Inches(0.95), border_color=clr)
    tb_t = add_text_box(s, Inches(1.1), y + Inches(0.05), Inches(4.0), Inches(0.35))
    first_paragraph(_tf(tb_t), title, size=16, bold=True, color=clr)
    tb_d = add_text_box(s, Inches(1.1), y + Inches(0.38), Inches(11.0), Inches(0.5))
    dtf = _tf(tb_d)
    dtf.word_wrap = True
    first_paragraph(dtf, desc, size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Toolchain & Ecosystem
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "工具链与生态系统", "从数据获取到可视化的完整管线")

tools = [
    ("数据获取", [
        "fetch_terrain.py — ETOPO1 真实地形下载 (9 个预定义区域)",
        "fetch_rivers.py — 河流走廊场生成 (程序化 / 文件)",
        "generate_genesis.py — HYDE 3.3 历史人口初始条件生成",
    ], ACCENT),
    ("仿真运行", [
        "politeia <config>.cfg — 主仿真程序",
        "mpirun -np N politeia ... — MPI 并行运行",
        "--restart checkpoint.bin — 弹性重启",
    ], ACCENT3),
    ("分析可视化", [
        "plot_worldmap.py — 世界地图可视化 (交互/静态/动画)",
        "plot_phase_diagram.py — 9 面板相图",
        "plot_polities.py — 政体分析 5 面板",
        "param_scan.py — 参数空间扫描",
    ], ACCENT4),
    ("标准算例", [
        "adam_eve — 最小种群 (2 粒子) 起始演化",
        "genesis_hyde_100k — HYDE 3.3 校准全球 10 万 agent",
        "terrain compare — 6 种地形对比实验",
    ], ACCENT2),
]

for i, (cat, items, clr) in enumerate(tools):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(2.0) + row * Inches(2.6)
    h = Inches(2.4)
    card = add_card(s, x, y, Inches(5.5), h, border_color=clr)
    tb_c = add_text_box(s, x + Inches(0.2), y + Inches(0.08), Inches(5.0), Inches(0.35))
    first_paragraph(_tf(tb_c), cat, size=17, bold=True, color=clr)
    for j, item in enumerate(items):
        tb_i = add_text_box(s, x + Inches(0.3), y + Inches(0.55) + Inches(j * 0.45), Inches(4.8), Inches(0.4))
        itf = _tf(tb_i)
        itf.word_wrap = True
        first_paragraph(itf, f"▸ {item}", size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Future Vision
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s)
title_slide_header(s, "愿景与展望", "从 10⁵ 到 10⁹ — 模拟整个人类文明史")

vision_items = [
    ("全球 80 亿粒子模拟",
     "Morton Z-curve SFC 已就绪；目标：在 HPC 集群上实现 10⁵→10⁹ 粒子规模",
     ACCENT),
    ("完整的人类文明史",
     "从 10,000 BCE 农业革命前夕 → 经过农业革命、工业革命 → 到现代社会",
     ACCENT3),
    ("可证伪的科学预测",
     "每条交互规则是一个可证伪假说；消融实验确定文明演化的关键驱动力",
     ACCENT4),
    ("历史验证",
     "HYDE 3.3 人口数据校准；Gini / 寿命 / 政体类型 / 技术水平与历史记录对比",
     ACCENT2),
    ("物理学视角的社会科学",
     "用相变理论理解社会制度突变，用统计力学解释不平等的必然性",
     ORANGE),
]

for i, (title, desc, clr) in enumerate(vision_items):
    y = Inches(2.0) + Inches(i * 1.0)
    card = add_card(s, Inches(0.8), y, Inches(11.5), Inches(0.9), border_color=clr)
    tb_t = add_text_box(s, Inches(1.1), y + Inches(0.08), Inches(11.0), Inches(0.35))
    first_paragraph(_tf(tb_t), title, size=17, bold=True, color=clr)
    tb_d = add_text_box(s, Inches(1.1), y + Inches(0.45), Inches(11.0), Inches(0.4))
    dtf = _tf(tb_d)
    dtf.word_wrap = True
    first_paragraph(dtf, desc, size=13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, BG_DARK)

add_accent_line(s, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT)

tb = add_text_box(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.2))
tf = _tf(tb)
first_paragraph(tf, "Politeia", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_paragraph(tf, "基于 Langevin 动力学的人类文明演化模拟器", size=22, color=ACCENT, align=PP_ALIGN.CENTER)

tb = add_text_box(s, Inches(2.5), Inches(4.5), Inches(8.3), Inches(1.8))
tf = _tf(tb)
tf.word_wrap = True
first_paragraph(tf, "\u201c用物理学的方法，理解人类文明的涌现\u201d", size=20, color=ACCENT4, align=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=12)
add_paragraph(tf, "对称规则  ·  不对称结果  ·  涌现方法论", size=16, color=DIM, align=PP_ALIGN.CENTER)
add_paragraph(tf, "C++20 · MPI+OpenMP · Langevin-Jump Diffusion", size=14, color=DIM, align=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/home/wanwb/ONE/Politeia/Politeia_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
