#!/usr/bin/env python3
"""Generate Vibe Research x Politeia presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colour palette ──
BG_DARK    = RGBColor(0x0F, 0x0F, 0x1A)
BG_SECTION = RGBColor(0x12, 0x16, 0x2B)
BG_CARD    = RGBColor(0x1A, 0x22, 0x40)
BG_CARD_ALT= RGBColor(0x14, 0x1C, 0x36)
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_ORANGE= RGBColor(0xFB, 0x92, 0x3C)
ACCENT_RED   = RGBColor(0xF8, 0x71, 0x71)
ACCENT_YELLOW= RGBColor(0xFA, 0xCC, 0x15)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
DIM_GRAY   = RGBColor(0x64, 0x74, 0x8B)
DARK_LINE  = RGBColor(0x33, 0x3E, 0x5B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ─── helpers ───

def add_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _tf(shape):
    return shape.text_frame

def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def card(slide, l, t, w, h, fill=BG_CARD, border=None, radius=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def rect(slide, l, t, w, h, fill=ACCENT_BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def circle(slide, l, t, size, fill=ACCENT_BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def diamond(slide, l, t, size, fill=ACCENT_BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def arrow_right(slide, l, t, w, h, fill=DIM_GRAY):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def chevron(slide, l, t, w, h, fill=ACCENT_BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def p0(tf, text, sz=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, fn="Microsoft YaHei"):
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0); p.space_after = Pt(2)
    for r in p.runs:
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color; r.font.name = fn

def ap(tf, text, sz=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, fn="Microsoft YaHei",
       sb=Pt(4), sa=Pt(2)):
    p = tf.add_paragraph()
    p.text = text; p.alignment = align; p.space_before = sb; p.space_after = sa
    for r in p.runs:
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color; r.font.name = fn
    return p

def section_header(slide, title, subtitle=None):
    rect(slide, Inches(0.7), Inches(0.45), Inches(0.15), Inches(0.55), ACCENT_BLUE)
    t = tb(slide, Inches(1.0), Inches(0.4), Inches(10), Inches(0.7))
    tf = _tf(t); tf.word_wrap = True
    p0(tf, title, sz=30, bold=True, color=WHITE)
    if subtitle:
        ap(tf, subtitle, sz=14, color=MID_GRAY)

def number_badge(slide, l, t, num, color=ACCENT_BLUE, text_color=BG_DARK):
    c = circle(slide, l, t, Inches(0.45), fill=color)
    t2 = tb(slide, l, t + Inches(0.04), Inches(0.45), Inches(0.4))
    tf = _tf(t2)
    p0(tf, str(num), sz=18, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s, BG_DARK)

rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT_BLUE)

rect(s, Inches(1.2), Inches(2.1), Inches(4.5), Pt(3), ACCENT_BLUE)

t = tb(s, Inches(1.2), Inches(2.4), Inches(10.5), Inches(1.5))
tf = _tf(t); tf.word_wrap = True
p0(tf, "Vibe Research", sz=52, bold=True, color=WHITE)
ap(tf, "AI \u534f\u4f5c\u9a71\u52a8\u7684\u79d1\u5b66\u7814\u7a76\u65b0\u8303\u5f0f", sz=26, color=ACCENT_BLUE)

t = tb(s, Inches(1.2), Inches(4.2), Inches(10.5), Inches(1))
tf = _tf(t); tf.word_wrap = True
p0(tf, "\u4ee5 Politeia \u6587\u660e\u6f14\u5316\u6a21\u62df\u5668\u4e3a\u6848\u4f8b", sz=20, color=LIGHT_GRAY)
ap(tf, "Human-AI Collaboration \u00b7 Langevin Dynamics \u00b7 Civilization Simulation", sz=14, color=DIM_GRAY)

t = tb(s, Inches(1.2), Inches(6.2), Inches(4), Inches(0.4))
p0(_tf(t), "2026.05", sz=13, color=DIM_GRAY)

for i, (label, clr) in enumerate([("Vibe Research", ACCENT_BLUE), ("C++20 / MPI", ACCENT_GREEN),
                                    ("31 Phases", ACCENT_PURPLE), ("110+ Tests", ACCENT_ORANGE)]):
    x = Inches(7.5) + Inches(i * 1.45)
    card(s, x, Inches(5.5), Inches(1.3), Inches(0.7), fill=BG_CARD, border=clr)
    t2 = tb(s, x + Inches(0.1), Inches(5.55), Inches(1.1), Inches(0.6))
    tf2 = _tf(t2); tf2.word_wrap = True
    p0(tf2, label, sz=11, bold=True, color=clr, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 2 — What is Vibe Research?
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u4ec0\u4e48\u662f Vibe Research\uff1f", "\u4ece Vibe Coding \u5230 Vibe Research \u2014\u2014 \u79d1\u5b66\u7814\u7a76\u7684\u8303\u5f0f\u8f6c\u53d8")

quote_card = card(s, Inches(0.7), Inches(1.7), Inches(11.8), Inches(1.1), fill=RGBColor(0x1E, 0x29, 0x4F), border=ACCENT_PURPLE)
rect(s, Inches(0.7), Inches(1.7), Inches(0.12), Inches(1.1), ACCENT_PURPLE)
t = tb(s, Inches(1.1), Inches(1.82), Inches(11.2), Inches(0.9))
tf = _tf(t); tf.word_wrap = True
p0(tf, "\u201c\u7814\u7a76\u8005\u63d0\u4f9b\u9ad8\u5c42\u65b9\u5411\u3001\u521b\u9020\u6027\u76f4\u89c9\u548c\u6279\u5224\u6027\u8bc4\u4ef7\uff0cLLM Agent \u6267\u884c\u52b3\u52a8\u5bc6\u96c6\u578b\u4efb\u52a1\u2014\u2014\u6587\u732e\u7efc\u8ff0\u3001\u7f16\u7801\u5b9e\u73b0\u3001\u6570\u636e\u5206\u6790\u3001\u8bba\u6587\u64b0\u5199\u201d",
    sz=15, color=LIGHT_GRAY)
ap(tf, "\u2014\u2014 Feng & Liu, \u201cA Visionary Look at Vibe Researching\u201d (arXiv:2604.00945, 2026)", sz=12, color=DIM_GRAY)

spectrum_colors = [MID_GRAY, RGBColor(0x6B, 0x7B, 0x93), ACCENT_GREEN, ACCENT_BLUE, ACCENT_PURPLE]
spectrum_labels = ["\u4f20\u7edf\u7814\u7a76", "\u5de5\u5177\u8f85\u52a9", "AI for Science", "Vibe Research", "Auto Research"]
spectrum_descs  = ["\u5168\u4eba\u5de5", "\u8ba1\u7b97\u5de5\u5177", "AI\u7b97\u7279\u5b9a\u95ee\u9898", "\u4eba\u5f15\u5bfc+AI\u6267\u884c", "\u5168AI\u81ea\u4e3b"]

for i in range(5):
    x = Inches(0.9) + Inches(i * 2.42)
    y = Inches(3.2)
    chevron(s, x, y, Inches(2.3), Inches(0.7), fill=spectrum_colors[i])
    t2 = tb(s, x + Inches(0.3), y + Inches(0.08), Inches(1.8), Inches(0.35))
    p0(_tf(t2), spectrum_labels[i], sz=12, bold=True, color=WHITE if i < 2 else BG_DARK, align=PP_ALIGN.CENTER)
    t3 = tb(s, x + Inches(0.2), y + Inches(0.4), Inches(2.0), Inches(0.3))
    p0(_tf(t3), spectrum_descs[i], sz=10, color=WHITE if i < 2 else BG_DARK, align=PP_ALIGN.CENTER)

rect(s, Inches(0.7), Inches(4.1), Inches(11.8), Pt(1), DIM_GRAY)

t = tb(s, Inches(0.7), Inches(4.15), Inches(11.8), Inches(0.4))
p0(_tf(t), "AI \u53c2\u4e0e\u5ea6\u9012\u589e  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2192", sz=11, color=DIM_GRAY, align=PP_ALIGN.CENTER)

comparisons = [
    ("vs AI for Science", "AI \u66ff\u4ee3\u7279\u5b9a\u8ba1\u7b97\u6b65\u9aa4\uff08\u5982 AlphaFold\uff09\uff0c\u4f46\u7814\u7a76\u6d41\u7a0b\u4e0d\u53d8", ACCENT_GREEN),
    ("vs Auto Research", "AI \u5b8c\u5168\u81ea\u4e3b\uff0c\u4eba\u7c7b\u4e0d\u53c2\u4e0e\u2014\u2014\u4f46\u8d28\u91cf\u548c\u53ef\u9760\u6027\u96be\u4ee5\u4fdd\u8bc1", ACCENT_PURPLE),
    ("\u7c7b\u6bd4 Vibe Coding", "\u5f00\u53d1\u8005\u8bf4\u610f\u56fe\uff0cAI\u5199\u4ee3\u7801 \u2192 \u7814\u7a76\u8005\u8bf4\u65b9\u5411\uff0cAI\u6267\u884c\u7814\u7a76", ACCENT_BLUE),
]
for i, (title, desc, clr) in enumerate(comparisons):
    x = Inches(0.7) + Inches(i * 4.0)
    y = Inches(4.8)
    c = card(s, x, y, Inches(3.8), Inches(1.0), fill=BG_CARD, border=clr)
    t2 = tb(s, x + Inches(0.2), y + Inches(0.08), Inches(3.4), Inches(0.35))
    p0(_tf(t2), title, sz=14, bold=True, color=clr)
    t3 = tb(s, x + Inches(0.2), y + Inches(0.42), Inches(3.4), Inches(0.55))
    tf3 = _tf(t3); tf3.word_wrap = True
    p0(tf3, desc, sz=11, color=LIGHT_GRAY)

karpathy_t = tb(s, Inches(0.7), Inches(6.2), Inches(11.8), Inches(0.5))
tf_k = _tf(karpathy_t); tf_k.word_wrap = True
p0(tf_k, "\u7075\u611f\u6765\u6e90\uff1aAndrej Karpathy (2025) \u63d0\u51fa Vibe Coding \u2192 \u81ea\u7136\u5ef6\u4f38\u5230\u79d1\u5b66\u7814\u7a76\u9886\u57df", sz=12, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 3 — Five Principles
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "Vibe Research \u4e94\u5927\u539f\u5219", "\u4ee5 Politeia \u9879\u76ee\u4e3a\u5b9e\u8df5\u6848\u4f8b")

principles = [
    ("\u4eba\u7c7b\u4f5c\u4e3a\u521b\u610f\u603b\u76d1", "Human as Creative Director",
     "\u9009\u62e9\u95ee\u9898\u3001\u5224\u65ad\u53d1\u73b0\u3001\u505a\u6218\u7565\u51b3\u7b56",
     "\u7814\u7a76\u8005\u8bbe\u8ba1\u7269\u7406\u6a21\u578b\u548c\u54f2\u5b66\u6846\u67b6", ACCENT_BLUE),
    ("\u81ea\u7136\u8bed\u8a00\u4e3a\u4e3b\u754c\u9762", "Natural Language Interface",
     "\u7528\u65e5\u5e38\u8bed\u8a00\u6307\u6325 Agent\uff0c\u800c\u975e\u7f16\u5199\u811a\u672c",
     "\u901a\u8fc7\u5bf9\u8bdd\u9a71\u52a8 31 \u4e2a\u5f00\u53d1\u9636\u6bb5", ACCENT_PURPLE),
    ("\u59d4\u6258\u4f46\u6709\u76d1\u7763", "Delegation with Oversight",
     "\u4efb\u52a1\u4ea4\u7ed9 Agent\uff0c\u4eba\u5ba1\u6838\u4ea7\u51fa",
     "\u7269\u7406\u516c\u5f0f\u548c\u6a21\u578b\u53c2\u6570\u9700\u4eba\u7c7b\u786e\u8ba4", ACCENT_GREEN),
    ("\u8fed\u4ee3\u7cbe\u70bc", "Iterative Refinement",
     "\u6307\u4ee4\u2192\u6267\u884c\u2192\u8bc4\u4ef7\u2192\u91cd\u5b9a\u5411\uff0c\u5faa\u73af\u5f80\u590d",
     "\u6bcf\u4e2a Phase \u90fd\u7ecf\u8fc7\u591a\u8f6e\u8fed\u4ee3\u4f18\u5316", ACCENT_ORANGE),
    ("\u4eba\u7c7b\u62c5\u8d23", "Human Accountability",
     "\u7814\u7a76\u8005\u4e3a\u6bcf\u4e2a\u7ed3\u8bba\u8d1f\u8d23\uff0cAI \u662f\u5de5\u5177\u800c\u975e\u5171\u540c\u4f5c\u8005",
     "research-proposal.md \u662f\u552f\u4e00\u6743\u5a01\u6e90", ACCENT_RED),
]

for i, (zh_name, en_name, desc, politeia, clr) in enumerate(principles):
    y = Inches(1.7) + Inches(i * 1.1)
    c = card(s, Inches(0.7), y, Inches(11.8), Inches(0.95), fill=BG_CARD, border=clr)
    number_badge(s, Inches(1.0), y + Inches(0.22), i+1, color=clr, text_color=BG_DARK)
    t_name = tb(s, Inches(1.6), y + Inches(0.08), Inches(3.5), Inches(0.4))
    p0(_tf(t_name), zh_name, sz=17, bold=True, color=clr)
    t_en = tb(s, Inches(1.6), y + Inches(0.45), Inches(3.5), Inches(0.3))
    p0(_tf(t_en), en_name, sz=11, color=DIM_GRAY)
    t_desc = tb(s, Inches(5.3), y + Inches(0.12), Inches(3.5), Inches(0.4))
    tf_d = _tf(t_desc); tf_d.word_wrap = True
    p0(tf_d, desc, sz=13, color=LIGHT_GRAY)
    t_pol = tb(s, Inches(9.0), y + Inches(0.12), Inches(3.3), Inches(0.7))
    tf_p = _tf(t_pol); tf_p.word_wrap = True
    p0(tf_p, "\u25b8 " + politeia, sz=12, color=WHITE)


# ═══════════════════════════════════════════════════
# SLIDE 4 — Politeia Case Study
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "Politeia\uff1a\u4e00\u4e2a Vibe Research \u7684\u5b8c\u6574\u6848\u4f8b")

t = tb(s, Inches(0.7), Inches(1.5), Inches(11.8), Inches(0.5))
tf = _tf(t); tf.word_wrap = True
p0(tf, "\u7528 Langevin \u52a8\u529b\u5b66\u6a21\u62df\u4eba\u7c7b\u6587\u660e\u6f14\u5316 \u2014\u2014 \u5c06\u6bcf\u4e2a\u4eba\u5efa\u6a21\u4e3a\u201c\u7c92\u5b50\u201d\uff0c\u5728\u771f\u5b9e\u5730\u5f62\u4e0a\u6f14\u5316", sz=16, color=LIGHT_GRAY)

stats = [("15,000+", "\u884c C++ \u4ee3\u7801", ACCENT_BLUE),
         ("31", "\u5f00\u53d1\u9636\u6bb5", ACCENT_PURPLE),
         ("110+", "\u5355\u5143\u6d4b\u8bd5", ACCENT_GREEN),
         ("~44\u5468", "\u5f00\u53d1\u5468\u671f", ACCENT_ORANGE)]
for i, (num, label, clr) in enumerate(stats):
    x = Inches(0.7) + Inches(i * 3.1)
    c = card(s, x, Inches(2.2), Inches(2.8), Inches(1.0), fill=BG_CARD, border=clr)
    t2 = tb(s, x + Inches(0.2), Inches(2.28), Inches(1.5), Inches(0.5))
    p0(_tf(t2), num, sz=28, bold=True, color=clr)
    t3 = tb(s, x + Inches(1.6), Inches(2.35), Inches(1.1), Inches(0.4))
    p0(_tf(t3), label, sz=13, color=LIGHT_GRAY)

roles_header = [("\u89d2\u8272", Inches(1.8)), ("\u4eba\u7c7b\u7814\u7a76\u8005", Inches(4.5)), ("AI Agent", Inches(4.5))]
role_data = [
    ("\u6838\u5fc3\u7406\u8bba", "\u8bbe\u8ba1\u7269\u7406\u6a21\u578b\u3001\u54f2\u5b66\u63a8\u5bfc", "\u4ea4\u53c9\u68c0\u67e5\u4e00\u81f4\u6027"),
    ("\u65b9\u7a0b\u63a8\u5bfc", "\u9009\u62e9 Langevin \u6846\u67b6\u3001\u5b9a\u4e49\u81ea\u7531\u5ea6", "\u8f85\u52a9\u516c\u5f0f\u6392\u7248\u548c\u9a8c\u8bc1"),
    ("\u4ee3\u7801\u5b9e\u73b0", "\u5ba1\u6838\u67b6\u6784\u51b3\u7b56", "\u7f16\u5199 C++20 \u4ee3\u7801"),
    ("\u6d4b\u8bd5\u9a8c\u8bc1", "\u8bbe\u8ba1\u7269\u7406\u9a8c\u8bc1\u6807\u51c6", "\u5b9e\u73b0 110+ \u5355\u5143\u6d4b\u8bd5"),
    ("\u6587\u6863\u7ef4\u62a4", "\u7f16\u5199 research-proposal", "\u7ef4\u62a4 wiki\u3001CODE_GUIDE"),
    ("\u5b9e\u9a8c\u8fd0\u884c", "\u8bbe\u8ba1\u5b9e\u9a8c\u65b9\u6848", "\u53c2\u6570\u626b\u63cf\u3001\u7ed3\u679c\u5206\u6790"),
]

hdr_y = Inches(3.5)
card(s, Inches(0.7), hdr_y, Inches(11.8), Inches(0.45), fill=RGBColor(0x0D, 0x2E, 0x52))
x_offsets = [Inches(0.9), Inches(2.8), Inches(7.5)]
for j, (h_text, h_w) in enumerate(roles_header):
    t2 = tb(s, x_offsets[j], hdr_y + Inches(0.06), h_w, Inches(0.35))
    p0(_tf(t2), h_text, sz=13, bold=True, color=ACCENT_BLUE)

for i, (role, human, ai) in enumerate(role_data):
    y = hdr_y + Inches(0.5) + Inches(i * 0.52)
    bg = BG_CARD if i % 2 == 0 else BG_CARD_ALT
    card(s, Inches(0.7), y, Inches(11.8), Inches(0.46), fill=bg)
    t_r = tb(s, Inches(0.9), y + Inches(0.06), Inches(1.8), Inches(0.34))
    p0(_tf(t_r), role, sz=13, bold=True, color=WHITE)
    t_h = tb(s, Inches(2.8), y + Inches(0.06), Inches(4.5), Inches(0.34))
    p0(_tf(t_h), human, sz=12, color=LIGHT_GRAY)
    t_a = tb(s, Inches(7.5), y + Inches(0.06), Inches(4.5), Inches(0.34))
    p0(_tf(t_a), ai, sz=12, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════
# SLIDE 5 — LLM Wiki Pattern
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u77e5\u8bc6\u7ba1\u7406\uff1aLLM Wiki \u6a21\u5f0f", "AGENTS.md \u5b9a\u4e49\u7684\u4e09\u5c42\u6587\u6863\u67b6\u6784")

layers = [
    ("raw/", "\u7b2c\u4e00\u5c42\uff1a\u539f\u59cb\u8d44\u6599", "\u4e0d\u53ef\u53d8\uff0cAI \u53ea\u8bfb", "\u8bba\u6587 PDF \u00b7 HYDE 3.3 \u6570\u636e \u00b7 \u624b\u5199\u7b14\u8bb0",
     ACCENT_RED, BG_DARK),
    ("docs/", "\u7b2c\u4e8c\u5c42\uff1a\u6b63\u5f0f\u6280\u672f\u6587\u6863", "\u957f\u671f\u4fdd\u7559", "\u968f\u673a\u5206\u5e03 \u00b7 \u5e76\u884c\u8bbe\u8ba1 \u00b7 \u6784\u5efa\u8bf4\u660e",
     ACCENT_ORANGE, BG_DARK),
    ("wiki/", "\u7b2c\u4e09\u5c42\uff1aAI \u8f85\u52a9\u5bfc\u822a\u5c42", "AI \u7ef4\u62a4", "\u7d22\u5f15 \u00b7 \u65f6\u95f4\u7ebf \u00b7 \u6a21\u5757\u901f\u67e5 \u00b7 ADR",
     ACCENT_BLUE, BG_DARK),
]

for i, (path, name, attr, content, clr, _) in enumerate(layers):
    y = Inches(1.7) + Inches(i * 1.5)
    c = card(s, Inches(0.7), y, Inches(5.8), Inches(1.3), fill=BG_CARD, border=clr)
    rect(s, Inches(0.7), y, Inches(0.12), Inches(1.3), clr)
    t_p = tb(s, Inches(1.1), y + Inches(0.08), Inches(1.5), Inches(0.4))
    p0(_tf(t_p), path, sz=18, bold=True, color=clr, fn="Consolas")
    t_n = tb(s, Inches(2.7), y + Inches(0.08), Inches(3.5), Inches(0.4))
    p0(_tf(t_n), name, sz=14, bold=True, color=WHITE)
    t_a = tb(s, Inches(2.7), y + Inches(0.45), Inches(3.5), Inches(0.3))
    p0(_tf(t_a), attr, sz=12, color=DIM_GRAY)
    t_c = tb(s, Inches(1.1), y + Inches(0.8), Inches(5.2), Inches(0.4))
    tf_c = _tf(t_c); tf_c.word_wrap = True
    p0(tf_c, content, sz=12, color=LIGHT_GRAY)

ops_card = card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.8), fill=BG_CARD, border=ACCENT_PURPLE)
t_ops = tb(s, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_ops), "\u64cd\u4f5c\u6d41\u7a0b", sz=18, bold=True, color=ACCENT_PURPLE)

ops = [("Ingest", "\u6444\u5165\u65b0\u8d44\u6599 \u2192 \u66f4\u65b0 wiki", ACCENT_GREEN),
       ("Change", "\u4ee3\u7801/\u6587\u6863\u53d8\u66f4 \u2192 \u66f4\u65b0\u6a21\u5757\u9875", ACCENT_BLUE),
       ("Decision", "\u8bbe\u8ba1\u51b3\u7b56 \u2192 ADR \u8bb0\u5f55", ACCENT_ORANGE),
       ("Query", "\u6709\u4ef7\u503c\u7684\u5206\u6790 \u2192 \u5b58\u4e3a wiki \u9875", ACCENT_PURPLE),
       ("Lint", "\u5065\u5eb7\u68c0\u67e5 \u2192 \u4e00\u81f4\u6027\u9a8c\u8bc1", ACCENT_RED)]
for i, (op, desc, clr) in enumerate(ops):
    y = Inches(2.35) + Inches(i * 0.42)
    number_badge(s, Inches(7.1), y, i+1, clr, BG_DARK)
    t_op = tb(s, Inches(7.65), y + Inches(0.02), Inches(1.2), Inches(0.35))
    p0(_tf(t_op), op, sz=13, bold=True, color=clr)
    t_od = tb(s, Inches(8.9), y + Inches(0.02), Inches(3.3), Inches(0.35))
    p0(_tf(t_od), desc, sz=12, color=LIGHT_GRAY)

core_card = card(s, Inches(6.8), Inches(4.8), Inches(5.7), Inches(1.6), fill=RGBColor(0x1E, 0x29, 0x4F), border=ACCENT_YELLOW)
diamond(s, Inches(7.1), Inches(5.0), Inches(0.45), ACCENT_YELLOW)
t_core = tb(s, Inches(7.7), Inches(5.0), Inches(4.5), Inches(0.4))
p0(_tf(t_core), "\u6838\u5fc3\u7ea6\u675f", sz=16, bold=True, color=ACCENT_YELLOW)
t_core2 = tb(s, Inches(7.3), Inches(5.5), Inches(5.0), Inches(0.8))
tf_c2 = _tf(t_core2); tf_c2.word_wrap = True
p0(tf_c2, "research-proposal.md \u662f\u201c\u9879\u76ee\u7075\u9b42\u201d\u2014\u2014\u7269\u7406\u516c\u5f0f\u548c\u6a21\u578b\u53c2\u6570\u7684\u5b9e\u8d28\u6027\u53d8\u66f4\u9700\u8981\u4eba\u7c7b\u786e\u8ba4", sz=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 6 — Physics Model
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u6838\u5fc3\u7269\u7406\u6a21\u578b\uff1aLangevin-\u8df3\u8dc3\u6269\u6563\u65b9\u7a0b", "m\u00b7q\u0308 = \u2212\u2207V \u2212 \u03b3\u00b7q\u0307 + \u03be(t) + J(t)")

eq_card = card(s, Inches(0.7), Inches(1.7), Inches(11.8), Inches(0.8), fill=RGBColor(0x0D, 0x2E, 0x52), border=ACCENT_BLUE)
t_eq = tb(s, Inches(1.0), Inches(1.8), Inches(11.2), Inches(0.6))
tf_eq = _tf(t_eq); tf_eq.word_wrap = True
p0(tf_eq, "m\u00b7q\u0308 = \u2212\u2207V(q) \u2212 \u03b3\u00b7q\u0307 + \u03be(t) + J(t)", sz=26, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, fn="Consolas")

terms = [
    ("\u2212\u2207V", "\u786e\u5b9a\u6027\u529b", "\u5730\u5f62\u52bf + \u4eba\u9645 LJ \u52bf", "force/", "\u805a\u96c6\u3001\u8fc1\u5f99", ACCENT_BLUE),
    ("\u2212\u03b3\u00b7q\u0307", "\u8017\u6563\u529b", "\u5236\u5ea6\u8150\u5316\u3001\u77e5\u8bc6\u9057\u5fd8", "integrator/", "\u7cfb\u7edf\u4e0d\u53ef\u9006", ACCENT_RED),
    ("\u03be(t)", "\u9ad8\u65af\u767d\u566a\u58f0", "\u65e5\u5e38\u968f\u673a\u4e8b\u4ef6", "integrator/", "\u5e03\u6717\u8fd0\u52a8", ACCENT_GREEN),
    ("J(t)", "Poisson \u8df3\u8dc3", "\u6280\u672f\u7a81\u7834\u3001\u7600\u75ab", "interaction/", "\u975e\u8fde\u7eed\u8dc3\u8fc1", ACCENT_ORANGE),
]

for i, (sym, name, social, module, effect, clr) in enumerate(terms):
    y = Inches(2.8) + Inches(i * 1.05)
    c = card(s, Inches(0.7), y, Inches(5.5), Inches(0.9), fill=BG_CARD, border=clr)
    t_s = tb(s, Inches(1.0), y + Inches(0.08), Inches(1.2), Inches(0.4))
    p0(_tf(t_s), sym, sz=20, bold=True, color=clr, fn="Consolas")
    t_n = tb(s, Inches(2.3), y + Inches(0.08), Inches(2.0), Inches(0.4))
    p0(_tf(t_n), name, sz=14, bold=True, color=WHITE)
    t_so = tb(s, Inches(2.3), y + Inches(0.45), Inches(3.7), Inches(0.4))
    p0(_tf(t_so), social, sz=12, color=LIGHT_GRAY)

mapping_card = card(s, Inches(6.5), Inches(2.8), Inches(6.0), Inches(4.2), fill=BG_CARD, border=ACCENT_PURPLE)
t_m = tb(s, Inches(6.8), Inches(2.9), Inches(5.5), Inches(0.4))
p0(_tf(t_m), "\u5206\u5b50\u52a8\u529b\u5b66 \u2192 \u793e\u4f1a\u6620\u5c04", sz=16, bold=True, color=ACCENT_PURPLE)

mapping = [
    ("\u7c92\u5b50", "\u4e2a\u4f53\u4eba\u7c7b"),
    ("\u4f4d\u7f6e & \u901f\u5ea6", "\u5730\u7406\u4f4d\u7f6e + \u8fc1\u5f99"),
    ("\u52bf\u80fd V(x)", "\u5730\u5f62\u00b7\u8d44\u6e90\u00b7\u4eba\u9645\u529b"),
    ("\u6e29\u5ea6 T", "\u793e\u4f1a\u52a8\u8361\u7a0b\u5ea6"),
    ("\u76f8\u53d8", "\u5236\u5ea6\u7a81\u53d8\uff08\u9769\u547d\uff09"),
    ("\u51dd\u805a/\u6210\u6838", "\u805a\u96c6\u4e3a\u90e8\u843d\u00b7\u56fd\u5bb6"),
    ("Poisson \u8df3\u8dc3", "\u6280\u672f\u7a81\u7834\u00b7\u7600\u75ab"),
]
for i, (phys, soc) in enumerate(mapping):
    y = Inches(3.4) + Inches(i * 0.48)
    bg = BG_CARD_ALT if i % 2 == 0 else BG_CARD
    t_ph = tb(s, Inches(6.9), y, Inches(2.2), Inches(0.35))
    p0(_tf(t_ph), phys, sz=12, color=ACCENT_PURPLE)
    arrow_right(s, Inches(9.1), y + Inches(0.08), Inches(0.35), Inches(0.18), DIM_GRAY)
    t_so = tb(s, Inches(9.6), y, Inches(2.6), Inches(0.35))
    p0(_tf(t_so), soc, sz=12, color=WHITE)


# ═══════════════════════════════════════════════════
# SLIDE 7 — Emergence
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u6d8c\u73b0\u65b9\u6cd5\u8bba\uff1a\u5bf9\u79f0\u89c4\u5219 \u2192 \u4e0d\u5bf9\u79f0\u7ed3\u679c")

mechanisms = [
    ("\u8d44\u6e90\u4ea4\u6362", "\u0394w \u221d (A_i\u2212A_j)/(A_i+A_j)", "Gini: 0 \u2192 0.89", ACCENT_BLUE),
    ("\u6587\u5316\u4f20\u64ad", "\u76f8\u4f3c\u2192\u8d8b\u540c\uff0c\u4e0d\u540c\u2192\u6392\u65a5", "\u6587\u5316\u5708\u81ea\u53d1\u5f62\u6210", ACCENT_PURPLE),
    ("\u6280\u672f\u6f14\u5316 \u03b5", "\u6f02\u79fb + \u4f20\u64ad + Poisson \u8df3\u8dc3", "\u519c\u4e1a\u2192\u5de5\u4e1a\u9769\u547d", ACCENT_GREEN),
    ("\u5fe0\u8bda\u5ea6 L", "dL/dt = \u4fdd\u62a4\u2212\u7a0e\u6536\u2212\u6587\u5316\u8ddd\u79bb", "Band \u2192 Empire", ACCENT_ORANGE),
    ("\u7600\u75ab", "\u514d\u75ab\u5411\u91cf + SIR + \u5dee\u5f02\u5316", "\u54e5\u4f26\u5e03\u4ea4\u6362\u5f0f\u707e\u96be", ACCENT_RED),
]

for i, (name, rule, result, clr) in enumerate(mechanisms):
    col = i % 2
    row = i // 2
    if i == 4:
        col = 0; row = 2
    x = Inches(0.7) + col * Inches(6.15)
    y = Inches(1.7) + row * Inches(1.55)
    w = Inches(5.8) if i < 4 else Inches(11.8)
    c = card(s, x, y, w, Inches(1.35), fill=BG_CARD, border=clr)
    circle(s, x + Inches(0.15), y + Inches(0.12), Inches(0.35), clr)
    t_num = tb(s, x + Inches(0.15), y + Inches(0.15), Inches(0.35), Inches(0.3))
    p0(_tf(t_num), str(i+1), sz=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    t_n = tb(s, x + Inches(0.65), y + Inches(0.08), Inches(3.0), Inches(0.35))
    p0(_tf(t_n), name, sz=16, bold=True, color=clr)
    t_r = tb(s, x + Inches(0.65), y + Inches(0.48), w - Inches(1.0), Inches(0.35))
    p0(_tf(t_r), "\u89c4\u5219\uff1a" + rule, sz=12, color=LIGHT_GRAY)
    t_res = tb(s, x + Inches(0.65), y + Inches(0.85), w - Inches(1.0), Inches(0.35))
    p0(_tf(t_res), "\u2192 \u6d8c\u73b0\uff1a" + result, sz=13, bold=True, color=WHITE)

t_motto = tb(s, Inches(0.7), Inches(6.6), Inches(11.8), Inches(0.5))
tf_m = _tf(t_motto); tf_m.word_wrap = True
p0(tf_m, "\u6838\u5fc3\u601d\u60f3\uff1a\u4e0d\u9884\u8bbe\u4efb\u4f55\u793e\u4f1a\u7ed3\u6784\uff0c\u8ba9\u4e00\u5207\u4ece\u5bf9\u79f0\u7684\u5fae\u89c2\u89c4\u5219\u4e2d\u81ea\u7136\u6d8c\u73b0", sz=15, color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 8 — Scientific Findings
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u79d1\u5b66\u53d1\u73b0\uff1a\u6d8c\u73b0\u9a8c\u8bc1\u7ed3\u679c")

metrics_card = card(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(3.0), fill=BG_CARD, border=ACCENT_BLUE)
t_mc = tb(s, Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_mc), "\u5e8f\u53c2\u91cf\u6d8c\u73b0\u503c (1000 \u7c92\u5b50)", sz=16, bold=True, color=ACCENT_BLUE)

metrics = [
    ("H", "\u5c42\u7ea7\u6df1\u5ea6", "5\u20147", "\u5b98\u50da\u5e1d\u56fd", ACCENT_BLUE),
    ("C", "\u72ec\u7acb\u653f\u4f53", "110\u2014254", "\u81ea\u53d1\u5f62\u6210\u201c\u56fd\u5bb6\u201d", ACCENT_PURPLE),
    ("Gini", "\u8d22\u5bcc\u4e0d\u5e73\u7b49", "~0.89", "\u5bf9\u79f0\u89c4\u5219\u2192\u6781\u5ea6\u4e0d\u5e73\u7b49", ACCENT_RED),
    ("\u27e8L\u27e9", "\u5fe0\u8bda\u5ea6", "0.78\u20140.94", "\u4e16\u88ad\u4f7f\u6743\u529b\u66f4\u7a33\u5b9a", ACCENT_GREEN),
    ("\u03a8", "\u5c01\u5efa\u5316", "0.52\u20140.69", "\u4e2d\u7b49\u5c01\u5efa\u7ed3\u6784", ACCENT_ORANGE),
]
for i, (sym, name, val, desc, clr) in enumerate(metrics):
    y = Inches(2.35) + Inches(i * 0.46)
    t_sym = tb(s, Inches(1.0), y, Inches(0.7), Inches(0.35))
    p0(_tf(t_sym), sym, sz=14, bold=True, color=clr, fn="Consolas")
    t_name = tb(s, Inches(1.7), y, Inches(1.5), Inches(0.35))
    p0(_tf(t_name), name, sz=12, color=LIGHT_GRAY)
    t_val = tb(s, Inches(3.3), y, Inches(1.2), Inches(0.35))
    p0(_tf(t_val), val, sz=14, bold=True, color=clr, align=PP_ALIGN.CENTER)
    t_desc = tb(s, Inches(4.5), y, Inches(1.8), Inches(0.35))
    p0(_tf(t_desc), desc, sz=11, color=DIM_GRAY)

geo_card = card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(3.0), fill=BG_CARD, border=ACCENT_ORANGE)
t_gc = tb(s, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_gc), "\u5730\u5f62\u5bf9\u6bd4\u5b9e\u9a8c", sz=16, bold=True, color=ACCENT_ORANGE)

geo_data = [
    ("\u4e2d\u56fd (\u5927\u5e73\u539f)", "H=7", "15\u4e2a", "0.088", ACCENT_RED),
    ("\u6b27\u6d32 (\u5c71\u8109\u5206\u5272)", "H=5", "27\u4e2a", "0.052", ACCENT_BLUE),
    ("\u6cb3\u8c37 (River)", "H=5", "\u2014", "0.080", ACCENT_GREEN),
    ("\u9ad8\u65af\u5355\u4e2d\u5fc3", "H=3", "30\u4e2a", "\u2014", ACCENT_PURPLE),
]
hdr_g = Inches(2.3)
for j, h in enumerate(["\u5730\u5f62", "H", "\u653f\u4f53", "HHI"]):
    x = Inches(7.2) + Inches(j * 1.35)
    t_h = tb(s, x, hdr_g, Inches(1.3), Inches(0.3))
    p0(_tf(t_h), h, sz=11, bold=True, color=ACCENT_ORANGE)

for i, (terrain, h, pol, hhi, clr) in enumerate(geo_data):
    y = hdr_g + Inches(0.35) + Inches(i * 0.4)
    for j, val in enumerate([terrain, h, pol, hhi]):
        x = Inches(7.2) + Inches(j * 1.35)
        t_v = tb(s, x, y, Inches(1.3), Inches(0.3))
        p0(_tf(t_v), val, sz=12, color=WHITE if j == 0 else clr)

findings_cards = [
    ("\u4e0d\u5e73\u7b49\u81ea\u53d1\u6d8c\u73b0", "Gini: 0 \u2192 0.89\uff0c\u5e42\u5f8b\u5206\u5e03", ACCENT_RED),
    ("\u6587\u5316\u5708\u81ea\u53d1\u5f62\u6210", "\u65e0\u9700\u9884\u8bbe\u79bb\u6563\u6807\u7b7e", ACCENT_PURPLE),
    ("\u653f\u6cbb\u5236\u5ea6\u6d8c\u73b0", "Band \u2192 Tribe \u2192 State \u2192 Empire", ACCENT_GREEN),
    ("\u5730\u7406\u51b3\u5b9a\u653f\u6cbb\u7ed3\u6784", "\u5e73\u539f\u2192\u96c6\u6743\uff0c\u5c71\u8109\u2192\u788e\u88c2", ACCENT_ORANGE),
    ("\u5a5a\u914d\u5236\u5ea6\u6d8c\u73b0", "\u8d2b\u7a77\u2192\u4e00\u592b\u4e00\u59bb\uff0c\u5bcc\u88d5\u2192\u4e00\u592b\u591a\u59bb", ACCENT_YELLOW),
]
for i, (title, desc, clr) in enumerate(findings_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.7) + col * Inches(4.1)
    y = Inches(5.0) + row * Inches(1.1)
    w = Inches(3.8) if i < 3 else Inches(3.8)
    c = card(s, x, y, w, Inches(0.95), fill=BG_CARD, border=clr)
    t_t = tb(s, x + Inches(0.2), y + Inches(0.05), w - Inches(0.3), Inches(0.35))
    p0(_tf(t_t), title, sz=14, bold=True, color=clr)
    t_d = tb(s, x + Inches(0.2), y + Inches(0.42), w - Inches(0.3), Inches(0.45))
    tf_d = _tf(t_d); tf_d.word_wrap = True
    p0(tf_d, desc, sz=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 9 — Technical Architecture
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u6280\u672f\u67b6\u6784\uff1aAI \u534f\u4f5c\u6784\u5efa\u7684 HPC \u7cfb\u7edf", "C++20 \u00b7 SoA \u00b7 MPI+OpenMP \u00b7 Morton Z-curve")

tech_items = [
    ("C++20", "\u8bed\u8a00\u6807\u51c6", "std::ranges, concepts", ACCENT_BLUE),
    ("SoA \u5e03\u5c40", "\u5185\u5b58\u4f18\u5316", "\u8fde\u7eed\u6570\u7ec4\uff0cSIMD \u53cb\u597d", ACCENT_GREEN),
    ("BBK \u79ef\u5206\u5668", "\u65f6\u95f4\u6b65\u8fdb", "Langevin \u4e8c\u9636\u79ef\u5206", ACCENT_PURPLE),
    ("Cell List", "\u90bb\u5c45\u641c\u7d22", "O(N) \u6876\u6392\u5e8f", ACCENT_ORANGE),
]
for i, (name, cat, desc, clr) in enumerate(tech_items):
    x = Inches(0.7) + Inches(i * 3.1)
    c = card(s, x, Inches(1.7), Inches(2.8), Inches(1.2), fill=BG_CARD, border=clr)
    t_n = tb(s, x + Inches(0.15), Inches(1.78), Inches(2.5), Inches(0.4))
    p0(_tf(t_n), name, sz=16, bold=True, color=clr)
    t_c = tb(s, x + Inches(0.15), Inches(2.15), Inches(2.5), Inches(0.3))
    p0(_tf(t_c), cat, sz=11, color=DIM_GRAY)
    t_d = tb(s, x + Inches(0.15), Inches(2.45), Inches(2.5), Inches(0.35))
    p0(_tf(t_d), desc, sz=12, color=LIGHT_GRAY)

par_card = card(s, Inches(0.7), Inches(3.2), Inches(11.8), Inches(3.8), fill=BG_CARD, border=ACCENT_BLUE)
t_par = tb(s, Inches(1.0), Inches(3.3), Inches(5), Inches(0.4))
p0(_tf(t_par), "\u5e76\u884c\u67b6\u6784\uff1aMPI + OpenMP \u6df7\u5408", sz=18, bold=True, color=ACCENT_BLUE)

pipeline_steps = [
    ("MPI\u533a\u57df\u5206\u89e3", "Morton Z-curve\nSFC", ACCENT_BLUE),
    ("\u8d1f\u8f7d\u5747\u8861", "\u6548\u7387 97%+", ACCENT_GREEN),
    ("OpenMP\u7ebf\u7a0b", "\u529b\u8ba1\u7b97\u3001\u79ef\u5206\u5668", ACCENT_PURPLE),
    ("Checkpoint", "N\u2192M \u5f39\u6027\u6269\u7f29", ACCENT_ORANGE),
]
for i, (step, desc, clr) in enumerate(pipeline_steps):
    x = Inches(1.0) + Inches(i * 2.9)
    y = Inches(3.9)
    chevron(s, x, y, Inches(2.6), Inches(0.8), clr)
    t_s = tb(s, x + Inches(0.35), y + Inches(0.08), Inches(2.0), Inches(0.35))
    p0(_tf(t_s), step, sz=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    t_d = tb(s, x + Inches(0.35), y + Inches(0.4), Inches(2.0), Inches(0.35))
    p0(_tf(t_d), desc, sz=10, color=BG_DARK, align=PP_ALIGN.CENTER)

modules = [
    ("core/", "\u57fa\u7840\u6570\u636e\u5c42", ACCENT_BLUE), ("domain/", "\u7a7a\u95f4\u7ba1\u7406", ACCENT_GREEN),
    ("force/", "\u529b\u8ba1\u7b97", ACCENT_RED), ("integrator/", "\u65f6\u95f4\u79ef\u5206", ACCENT_PURPLE),
    ("interaction/", "\u4e2a\u4f53\u4ea4\u4e92", ACCENT_ORANGE), ("population/", "\u4eba\u53e3\u52a8\u529b\u5b66", ACCENT_YELLOW),
    ("analysis/", "\u5206\u6790\u89c2\u6d4b", ACCENT_GREEN), ("io/", "\u8f93\u5165\u8f93\u51fa", DIM_GRAY),
]
for i, (mod, desc, clr) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = Inches(1.0) + col * Inches(2.9)
    y = Inches(5.1) + row * Inches(0.85)
    c = card(s, x, y, Inches(2.6), Inches(0.7), fill=BG_CARD_ALT, border=clr)
    t_m = tb(s, x + Inches(0.1), y + Inches(0.05), Inches(1.6), Inches(0.35))
    p0(_tf(t_m), mod, sz=13, bold=True, color=clr, fn="Consolas")
    t_d = tb(s, x + Inches(0.1), y + Inches(0.35), Inches(2.3), Inches(0.3))
    p0(_tf(t_d), desc, sz=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 10 — Development Timeline
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u5f00\u53d1\u5386\u7a0b\uff1a31 \u4e2a Phase \u7684 Vibe Research \u5b9e\u8df5", "~44 \u5468 \u00b7 \u6d8c\u73b0\u4f18\u5148\u539f\u5219")

phases = [
    ("Phase 0\u20143", "\u57fa\u7840\u6784\u5efa", "\u9aa8\u67b6 + \u7a7a\u95f4 + \u8d44\u6e90\u4ea4\u6362 + Langevin", ACCENT_BLUE),
    ("Phase 4\u20146", "\u793e\u4f1a\u673a\u5236", "\u6587\u5316 + \u4eba\u53e3 + \u6280\u672f\u7a81\u7834", ACCENT_PURPLE),
    ("Phase 7\u201412", "\u5e76\u884c\u5316", "MPI + SFC + OpenMP + \u771f\u5b9e\u5730\u5f62", ACCENT_GREEN),
    ("Phase 13\u201416", "\u5c42\u7ea7\u5236\u5ea6", "\u5fe0\u8bda\u5ea6 + \u4e16\u88ad + \u653f\u4f53\u68c0\u6d4b", ACCENT_ORANGE),
    ("Phase 17\u201422", "\u5b9e\u9a8c\u9a8c\u8bc1", "Scaling + \u5730\u5f62\u5bf9\u6bd4 + \u76f8\u53d8\u68c0\u6d4b", ACCENT_RED),
    ("Phase 23\u201426", "\u6846\u67b6\u5b8c\u5584", "\u627f\u8f7d\u529b + \u4e2d\u56fdvs\u6b27\u6d32 + \u6d88\u878d\u5b9e\u9a8c", ACCENT_YELLOW),
    ("Phase 27\u201431", "\u751f\u6001\u5b8c\u5584", "\u4e16\u754c\u5730\u56fe + HYDE\u6821\u51c6 + Checkpoint", ACCENT_BLUE),
]

rect(s, Inches(2.0), Inches(1.9), Pt(3), Inches(5.0), DIM_GRAY)

for i, (phase, cat, desc, clr) in enumerate(phases):
    y = Inches(1.7) + Inches(i * 0.72)
    circle(s, Inches(1.88), y + Inches(0.12), Inches(0.28), clr)
    t_ph = tb(s, Inches(2.5), y + Inches(0.02), Inches(2.0), Inches(0.35))
    p0(_tf(t_ph), phase, sz=14, bold=True, color=clr)
    t_cat = tb(s, Inches(4.5), y + Inches(0.02), Inches(1.8), Inches(0.35))
    p0(_tf(t_cat), cat, sz=13, bold=True, color=WHITE)
    t_desc = tb(s, Inches(6.5), y + Inches(0.02), Inches(5.5), Inches(0.35))
    p0(_tf(t_desc), desc, sz=12, color=LIGHT_GRAY)
    t_check = tb(s, Inches(11.8), y + Inches(0.02), Inches(0.5), Inches(0.35))
    p0(_tf(t_check), "\u2713", sz=16, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

stat_items = [("15,000+", "\u884c\u4ee3\u7801"), ("110+", "\u5355\u5143\u6d4b\u8bd5"),
              ("31", "\u5f00\u53d1\u9636\u6bb5"), ("10\u2075", "\u7c92\u5b50\u9a8c\u8bc1")]
for i, (num, label) in enumerate(stat_items):
    x = Inches(0.7) + Inches(i * 3.1)
    c = card(s, x, Inches(6.8), Inches(2.8), Inches(0.55), fill=BG_CARD, border=ACCENT_BLUE)
    t_n = tb(s, x + Inches(0.2), Inches(6.84), Inches(1.2), Inches(0.4))
    p0(_tf(t_n), num, sz=18, bold=True, color=ACCENT_BLUE)
    t_l = tb(s, x + Inches(1.4), Inches(6.87), Inches(1.3), Inches(0.35))
    p0(_tf(t_l), label, sz=12, color=DIM_GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 11 — Vibe Research Practices
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "Vibe Research \u5728 Politeia \u4e2d\u7684\u5173\u952e\u5b9e\u8df5")

practices = [
    ("\u4eba\u7c7b\u628a\u63a7\u7269\u7406\u6743\u5a01",
     ["research-proposal.md \u662f\u552f\u4e00\u6743\u5a01\u6e90",
      "AI \u53ef\u8f85\u52a9\u7f16\u8f91\uff0c\u4f46\u65b9\u7a0b\u548c\u53c2\u6570\u9700\u4eba\u7c7b\u786e\u8ba4",
      "\u5f15\u7528\u6eaf\u6e90\uff1awiki \u9875\u9762\u5fc5\u987b\u6807\u6ce8\u51fa\u5904"],
     ACCENT_BLUE),
    ("AI \u9a71\u52a8\u7684\u77e5\u8bc6\u7ba1\u7406",
     ["AGENTS.md \u5b9a\u4e49 Ingest/Change/Decision/Query/Lint \u5de5\u4f5c\u6d41",
      "AI \u7ef4\u62a4 wiki \u7d22\u5f15\u3001\u65f6\u95f4\u7ebf\u3001\u6a21\u5757\u901f\u67e5",
      "\u5b9a\u671f Lint \u68c0\u67e5\u4e00\u81f4\u6027\u3001\u5b64\u7acb\u9875\u9762\u3001\u8fc7\u65f6\u58f0\u660e"],
     ACCENT_PURPLE),
    ("\u8fed\u4ee3\u5f00\u53d1\u5faa\u73af",
     ["\u4eba\u7c7b\u63d0\u51fa\u76ee\u6807 \u2192 AI \u5b9e\u73b0\u4ee3\u7801 \u2192 \u6d4b\u8bd5\u9a8c\u8bc1 \u2192 \u4eba\u7c7b\u5ba1\u6838",
      "31 \u4e2a Phase \u90fd\u9075\u5faa\u6b64\u5faa\u73af",
      "\u6bcf\u4e2a Phase \u591a\u8f6e\u8fed\u4ee3\u76f4\u81f3\u901a\u8fc7\u7269\u7406\u9a8c\u8bc1"],
     ACCENT_GREEN),
    ("\u53ef\u8bc1\u4f2a\u7684\u79d1\u5b66\u65b9\u6cd5",
     ["\u6bcf\u6761\u4ea4\u4e92\u89c4\u5219\u662f\u4e00\u4e2a\u53ef\u8bc1\u4f2a\u5047\u8bf4",
      "\u6d88\u878d\u5b9e\u9a8c\uff08\u6a21\u5757\u5f00\u5173\uff09\u786e\u5b9a\u5173\u952e\u9a71\u52a8\u529b",
      "\u5386\u53f2\u6570\u636e\u6821\u51c6\uff08HYDE 3.3\uff09\u9a8c\u8bc1\u6d8c\u73b0\u7ed3\u679c"],
     ACCENT_ORANGE),
]

for i, (title, items, clr) in enumerate(practices):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * Inches(6.15)
    y = Inches(1.7) + row * Inches(2.7)
    c = card(s, x, y, Inches(5.8), Inches(2.5), fill=BG_CARD, border=clr)
    number_badge(s, x + Inches(0.2), y + Inches(0.15), i+1, clr, BG_DARK)
    t_t = tb(s, x + Inches(0.8), y + Inches(0.12), Inches(4.5), Inches(0.4))
    p0(_tf(t_t), title, sz=16, bold=True, color=clr)
    for j, item in enumerate(items):
        t_i = tb(s, x + Inches(0.4), y + Inches(0.65) + Inches(j * 0.5), Inches(5.0), Inches(0.45))
        tf_i = _tf(t_i); tf_i.word_wrap = True
        p0(tf_i, "\u25b8 " + item, sz=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 12 — Advantages & Challenges
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "Vibe Research \u7684\u4f18\u52bf\u4e0e\u6311\u6218")

adv_card = card(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(5.0), fill=BG_CARD, border=ACCENT_GREEN)
t_adv = tb(s, Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_adv), "\u2713  \u4f18\u52bf\uff08\u5728 Politeia \u4e2d\u9a8c\u8bc1\uff09", sz=18, bold=True, color=ACCENT_GREEN)

advantages = [
    ("\u5f00\u53d1\u901f\u5ea6", "44 \u5468\u5b8c\u6210 31 Phase\uff0c\u4f20\u7edf\u65b9\u5f0f\u53ef\u80fd\u9700 2-3 \u5e74"),
    ("\u4ee3\u7801\u8d28\u91cf", "110+ \u5355\u5143\u6d4b\u8bd5\uff0c\u6301\u7eed\u96c6\u6210\u9a8c\u8bc1"),
    ("\u77e5\u8bc6\u7cfb\u7edf\u5316", "\u4e09\u5c42\u6587\u6863\u67b6\u6784 + \u81ea\u52a8\u7ef4\u62a4\u7684 wiki"),
    ("\u8de8\u9886\u57df\u6574\u5408", "\u7269\u7406\u5b66 + \u793e\u4f1a\u79d1\u5b66 + HPC + \u8003\u53e4\u5b66"),
    ("\u8ba4\u77e5\u8d1f\u8f7d\u964d\u4f4e", "\u7814\u7a76\u8005\u4e13\u6ce8\u521b\u9020\u6027\u601d\u8003\uff0cAI \u5904\u7406\u7e41\u91cd\u6267\u884c"),
]
for i, (title, desc) in enumerate(advantages):
    y = Inches(2.4) + Inches(i * 0.82)
    circle(s, Inches(1.1), y + Inches(0.05), Inches(0.3), ACCENT_GREEN)
    t_n = tb(s, Inches(1.1), y + Inches(0.07), Inches(0.3), Inches(0.25))
    p0(_tf(t_n), "\u2713", sz=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    t_t = tb(s, Inches(1.6), y, Inches(4.5), Inches(0.35))
    p0(_tf(t_t), title, sz=14, bold=True, color=WHITE)
    t_d = tb(s, Inches(1.6), y + Inches(0.35), Inches(4.5), Inches(0.4))
    tf_d = _tf(t_d); tf_d.word_wrap = True
    p0(tf_d, desc, sz=11, color=LIGHT_GRAY)

chal_card = card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0), fill=BG_CARD, border=ACCENT_RED)
t_chal = tb(s, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_chal), "\u26a0  \u6311\u6218\u4e0e\u5e94\u5bf9", sz=18, bold=True, color=ACCENT_RED)

challenges = [
    ("\u5e7b\u89c9\u98ce\u9669", "research-proposal \u4f5c\u4e3a\u552f\u4e00\u6743\u5a01\u6e90"),
    ("\u7269\u7406\u6b63\u786e\u6027", "\u4eba\u7c7b\u5ba1\u6838\u6240\u6709\u65b9\u7a0b\u548c\u5e38\u91cf"),
    ("\u53ef\u91cd\u590d\u6027", "\u6807\u51c6\u7b97\u4f8b\u4f53\u7cfb + Checkpoint/Restart"),
    ("\u5f52\u5c5e\u4e0e\u8d23\u4efb", "\u4eba\u7c7b\u5bf9\u6240\u6709\u79d1\u5b66\u7ed3\u8bba\u8d1f\u8d23"),
    ("\u8d8b\u540c\u601d\u7ef4\u98ce\u9669", "\u6d88\u878d\u5b9e\u9a8c + \u591a\u89d2\u5ea6\u9a8c\u8bc1"),
]
for i, (title, solution) in enumerate(challenges):
    y = Inches(2.4) + Inches(i * 0.82)
    circle(s, Inches(7.2), y + Inches(0.05), Inches(0.3), ACCENT_RED)
    t_n = tb(s, Inches(7.2), y + Inches(0.07), Inches(0.3), Inches(0.25))
    p0(_tf(t_n), "!", sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t_t = tb(s, Inches(7.7), y, Inches(4.5), Inches(0.35))
    p0(_tf(t_t), title, sz=14, bold=True, color=WHITE)
    t_s = tb(s, Inches(7.7), y + Inches(0.35), Inches(4.5), Inches(0.4))
    tf_s = _tf(t_s); tf_s.word_wrap = True
    p0(tf_s, "\u2192 " + solution, sz=11, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════
# SLIDE 13 — Toolchain
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u5de5\u5177\u94fe\u4e0e\u751f\u6001\u7cfb\u7edf", "\u4ece\u6570\u636e\u83b7\u53d6\u5230\u53ef\u89c6\u5316\u7684\u5b8c\u6574\u7ba1\u7ebf")

pipeline_phases = [
    ("\u6570\u636e\u83b7\u53d6", ["fetch_terrain.py", "fetch_rivers.py", "generate_genesis.py"],
     "ETOPO1 \u5730\u5f62\u00b7\u6cb3\u6d41\u00b7HYDE 3.3", ACCENT_BLUE),
    ("\u4eff\u771f\u8fd0\u884c", ["politeia <cfg>", "mpirun -np N ...", "--restart ckpt.bin"],
     "C++20 \u00b7 MPI+OpenMP", ACCENT_GREEN),
    ("\u5206\u6790\u53ef\u89c6\u5316", ["plot_worldmap.py", "plot_phase_diagram.py", "plot_polities.py"],
     "\u4ea4\u4e92\u5730\u56fe\u00b7\u76f8\u56fe\u00b7\u653f\u4f53\u5206\u6790", ACCENT_PURPLE),
]

for i, (phase, tools, desc, clr) in enumerate(pipeline_phases):
    x = Inches(0.7) + Inches(i * 4.15)
    c = card(s, x, Inches(1.7), Inches(3.85), Inches(3.5), fill=BG_CARD, border=clr)
    rect(s, x, Inches(1.7), Inches(3.85), Inches(0.55), clr)
    t_p = tb(s, x + Inches(0.15), Inches(1.78), Inches(3.5), Inches(0.4))
    p0(_tf(t_p), phase, sz=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    for j, tool in enumerate(tools):
        t_t = tb(s, x + Inches(0.3), Inches(2.5) + Inches(j * 0.55), Inches(3.3), Inches(0.45))
        tf_t = _tf(t_t); tf_t.word_wrap = True
        p0(tf_t, "\u25b8 " + tool, sz=12, color=WHITE, fn="Consolas")
    t_d = tb(s, x + Inches(0.2), Inches(4.5), Inches(3.5), Inches(0.5))
    tf_d = _tf(t_d); tf_d.word_wrap = True
    p0(tf_d, desc, sz=11, color=DIM_GRAY, align=PP_ALIGN.CENTER)

if i < 2:
    arrow_right(s, x + Inches(3.95), Inches(3.2), Inches(0.35), Inches(0.3), DIM_GRAY)

cases_card = card(s, Inches(0.7), Inches(5.5), Inches(11.8), Inches(1.5), fill=BG_CARD, border=ACCENT_ORANGE)
t_cc = tb(s, Inches(1.0), Inches(5.6), Inches(3.0), Inches(0.4))
p0(_tf(t_cc), "\u6807\u51c6\u7b97\u4f8b\u4f53\u7cfb", sz=16, bold=True, color=ACCENT_ORANGE)

cases = [
    ("adam_eve", "2 \u7c92\u5b50\u8d77\u59cb\u6f14\u5316", ACCENT_GREEN),
    ("genesis_hyde_100k", "HYDE 3.3 \u6821\u51c6 10\u4e07 agent", ACCENT_BLUE),
    ("terrain compare", "6 \u79cd\u5730\u5f62\u5bf9\u6bd4\u5b9e\u9a8c", ACCENT_PURPLE),
]
for i, (name, desc, clr) in enumerate(cases):
    x = Inches(1.0) + Inches(i * 3.9)
    y = Inches(6.15)
    t_n = tb(s, x, y, Inches(2.5), Inches(0.35))
    p0(_tf(t_n), name, sz=14, bold=True, color=clr, fn="Consolas")
    t_d = tb(s, x, y + Inches(0.35), Inches(3.5), Inches(0.3))
    p0(_tf(t_d), desc, sz=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# SLIDE 14 — Future
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
section_header(s, "\u5c55\u671b\uff1aVibe Research \u7684\u672a\u6765", "\u4ece Politeia \u5230\u66f4\u5e7f\u6cdb\u7684\u79d1\u5b66\u7814\u7a76")

politeia_future = [
    ("\u5168\u7403 80 \u4ebf\u7c92\u5b50\u6a21\u62df", "Morton Z-curve SFC \u5df2\u5c31\u7eea\uff0c\u76ee\u6807 HPC \u96c6\u7fa4 10\u2079", ACCENT_BLUE),
    ("\u5b8c\u6574\u4eba\u7c7b\u6587\u660e\u53f2", "10,000 BCE \u2192 \u519c\u4e1a\u9769\u547d \u2192 \u5de5\u4e1a\u9769\u547d \u2192 \u73b0\u4ee3", ACCENT_GREEN),
    ("\u5386\u53f2\u6570\u636e\u6821\u51c6", "Gini/\u5bff\u547d/\u653f\u4f53 vs \u5386\u53f2\u8bb0\u5f55\u5bf9\u6bd4", ACCENT_ORANGE),
]

vr_future = [
    ("\u7814\u7a76\u534f\u4f5c\u5de5\u5177\u6210\u719f", "Deep Research, Agent Laboratory \u7b49\u5e73\u53f0\u5d1b\u8d77", ACCENT_PURPLE),
    ("\u591a Agent \u534f\u4f5c\u67b6\u6784", "\u6587\u732eAgent + \u5b9e\u9a8cAgent + \u5199\u4f5cAgent \u534f\u540c\u5de5\u4f5c", ACCENT_BLUE),
    ("\u53ef\u5ba1\u8ba1\u7684\u7814\u7a76\u8fc7\u7a0b", "\u4ea4\u4e92\u65e5\u5fd7 + \u7248\u672c\u63a7\u5236 + \u53ef\u590d\u73b0\u7ba1\u7ebf", ACCENT_GREEN),
    ("\u5168\u9762\u6e17\u900f", "\u4ece\u8f6f\u4ef6\u5de5\u7a0b\u5230\u81ea\u7136\u79d1\u5b66\u5230\u793e\u4f1a\u79d1\u5b66", ACCENT_ORANGE),
]

lc = card(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(3.3), fill=BG_CARD, border=ACCENT_BLUE)
t_lc = tb(s, Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_lc), "Politeia \u7684\u4e0b\u4e00\u6b65", sz=16, bold=True, color=ACCENT_BLUE)

for i, (title, desc, clr) in enumerate(politeia_future):
    y = Inches(2.4) + Inches(i * 0.85)
    number_badge(s, Inches(1.0), y, i+1, clr, BG_DARK)
    t_t = tb(s, Inches(1.6), y - Inches(0.02), Inches(4.5), Inches(0.35))
    p0(_tf(t_t), title, sz=14, bold=True, color=clr)
    t_d = tb(s, Inches(1.6), y + Inches(0.32), Inches(4.5), Inches(0.45))
    tf_d = _tf(t_d); tf_d.word_wrap = True
    p0(tf_d, desc, sz=11, color=LIGHT_GRAY)

rc = card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(3.3), fill=BG_CARD, border=ACCENT_PURPLE)
t_rc = tb(s, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p0(_tf(t_rc), "Vibe Research \u7684\u8d8b\u52bf", sz=16, bold=True, color=ACCENT_PURPLE)

for i, (title, desc, clr) in enumerate(vr_future):
    y = Inches(2.35) + Inches(i * 0.7)
    number_badge(s, Inches(7.1), y, i+1, clr, BG_DARK)
    t_t = tb(s, Inches(7.7), y - Inches(0.02), Inches(4.5), Inches(0.35))
    p0(_tf(t_t), title, sz=13, bold=True, color=clr)
    t_d = tb(s, Inches(7.7), y + Inches(0.28), Inches(4.5), Inches(0.35))
    tf_d = _tf(t_d); tf_d.word_wrap = True
    p0(tf_d, desc, sz=11, color=LIGHT_GRAY)

belief_card = card(s, Inches(0.7), Inches(5.3), Inches(11.8), Inches(1.6), fill=RGBColor(0x1E, 0x29, 0x4F), border=ACCENT_YELLOW)
t_b = tb(s, Inches(1.0), Inches(5.5), Inches(11.2), Inches(1.2))
tf_b = _tf(t_b); tf_b.word_wrap = True
p0(tf_b, "\u6838\u5fc3\u4fe1\u5ff5", sz=14, bold=True, color=ACCENT_YELLOW)
ap(tf_b, "\u6700\u597d\u7684\u79d1\u5b66\u7814\u7a76\u5c06\u6765\u81ea\u4eba\u7c7b\u521b\u9020\u529b\u4e0e AI \u6267\u884c\u529b\u7684\u6df1\u5ea6\u7ed3\u5408 \u2014\u2014 \u4eba\u7c7b\u63d0\u4f9b\u201c\u4e3a\u4ec0\u4e48\u201d\uff0cAI \u63d0\u4f9b\u201c\u600e\u4e48\u505a\u201d\u3002",
    sz=16, color=WHITE)
ap(tf_b, "Politeia \u8bc1\u660e\u4e86\u8fd9\u79cd\u6a21\u5f0f\u53ef\u4ee5\u4ea7\u51fa\u4e25\u8083\u7684\u3001\u53ef\u9a8c\u8bc1\u7684\u79d1\u5b66\u7814\u7a76\u6210\u679c\u3002",
    sz=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════
# SLIDE 15 — End
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s, BG_DARK)

rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT_BLUE)
rect(s, Inches(4.0), Inches(2.2), Inches(5.3), Pt(3), ACCENT_BLUE)

t = tb(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0))
tf = _tf(t)
p0(tf, "Vibe Research \u00d7 Politeia", sz=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

t = tb(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6))
tf = _tf(t)
p0(tf, "\u8ba9\u4eba\u7c7b\u4e13\u6ce8\u4e8e\u521b\u9020\u6027\u601d\u8003\uff0c\u8ba9 AI \u5904\u7406\u7e41\u91cd\u7684\u6267\u884c", sz=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

t = tb(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.0))
tf = _tf(t); tf.word_wrap = True
p0(tf, "Human as Creative Director  \u00b7  AI as Tireless Executor", sz=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
ap(tf, "", sz=8)
ap(tf, "\u5bf9\u79f0\u89c4\u5219  \u00b7  \u4e0d\u5bf9\u79f0\u7ed3\u679c  \u00b7  \u6d8c\u73b0\u65b9\u6cd5\u8bba", sz=14, color=DIM_GRAY, align=PP_ALIGN.CENTER)

ref_card = card(s, Inches(2.5), Inches(5.8), Inches(8.3), Inches(1.0), fill=BG_CARD)
t_ref = tb(s, Inches(2.7), Inches(5.9), Inches(7.9), Inches(0.8))
tf_ref = _tf(t_ref); tf_ref.word_wrap = True
p0(tf_ref, "\u53c2\u8003\u6587\u732e", sz=11, bold=True, color=DIM_GRAY)
ap(tf_ref, "Feng & Liu (2026), \u201cA Visionary Look at Vibe Researching\u201d, arXiv:2604.00945", sz=10, color=DIM_GRAY)
ap(tf_ref, "Karpathy (2025), \u201cVibe Coding\u201d", sz=10, color=DIM_GRAY)


# ── Save ──
output = "/home/wanwb/ONE/Politeia/Vibe_Research_Politeia.pptx"
prs.save(output)
print(f"Saved to {output}")
